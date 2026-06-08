#!/usr/bin/env python3
"""
Summer 2026 Camp Health & Safety — Counselor Training Deck
Generated from the Summer 2026 Staff Handbook (Section 9).
Audience: camp counselors (18-22). Tone: professional, concise, mild clean humor.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette ----------
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)   # headers / titles
TEAL   = RGBColor(0x2A, 0x9D, 0x8F)   # accent
RED    = RGBColor(0xE6, 0x3A, 0x46)   # emergency / alert
AMBER  = RGBColor(0xF4, 0xA2, 0x61)   # caution
LIGHT  = RGBColor(0xEE, 0xF2, 0xF6)   # placeholder fill
GRAY   = RGBColor(0x5A, 0x63, 0x6E)   # body / footer
DKGRAY = RGBColor(0x33, 0x3A, 0x42)   # body text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PHBORD = RGBColor(0xB8, 0xC4, 0xD0)   # placeholder border

FONT = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_rect(slide, x, y, w, h, color, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, lines, size=18, color=DKGRAY, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
        font=FONT, space_after=6, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        # allow per-line dict overrides
        if isinstance(ln, dict):
            text = ln.get("t", "")
            run = p.add_run(); run.text = text
            f = run.font
            f.size = Pt(ln.get("size", size)); f.bold = ln.get("bold", bold)
            f.italic = ln.get("italic", italic); f.name = font
            f.color.rgb = ln.get("color", color)
            if ln.get("bullet"):
                _bullet(p, ln.get("bcolor", TEAL))
            if ln.get("indent"):
                p.level = ln["indent"]
        else:
            run = p.add_run(); run.text = ln
            f = run.font
            f.size = Pt(size); f.bold = bold; f.italic = italic
            f.name = font; f.color.rgb = color
    return tb


def _bullet(p, color):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
    buClr = pPr.makeelement(qn('a:buClr'), {})
    srgb = pPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0], color[1], color[2])})
    buClr.append(srgb)
    pPr.append(buClr)
    pPr.append(buChar)


def bullets(slide, x, y, w, h, items, size=17, color=DKGRAY,
            bcolor=TEAL, space_after=9, lead=None):
    """items: list of strings OR (text, level) OR dict."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if lead:
        p = tf.paragraphs[0]; first = False
        r = p.add_run(); r.text = lead
        r.font.size = Pt(size + 1); r.font.bold = True
        r.font.name = FONT; r.font.color.rgb = NAVY
        p.space_after = Pt(space_after + 2)
    for it in items:
        level = 0
        if isinstance(it, tuple):
            text, level = it
        else:
            text = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        p.line_spacing = 1.0
        p.level = level
        r = p.add_run(); r.text = text
        r.font.size = Pt(size if level == 0 else size - 1)
        r.font.name = FONT
        r.font.color.rgb = color if level == 0 else GRAY
        r.font.bold = False
        _bullet(p, bcolor if level == 0 else AMBER)
    return tb


def header(slide, kicker, title, bar=NAVY, kicker_color=TEAL):
    # top accent bar
    fill_rect(slide, 0, 0, SW, Inches(1.25), bar)
    fill_rect(slide, 0, Inches(1.25), SW, Inches(0.06), TEAL)
    txt(slide, Inches(0.6), Inches(0.16), Inches(11), Inches(0.35),
        kicker.upper(), size=13, color=RGBColor(0xCB, 0xE3, 0xDF), bold=True)
    txt(slide, Inches(0.6), Inches(0.46), Inches(12), Inches(0.7),
        title, size=30, color=WHITE, bold=True)


def footer(slide, n):
    txt(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
        "Summer 2026 Camp Health & Safety  •  Counselor Training",
        size=10, color=GRAY)
    txt(slide, Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.35),
        str(n), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def placeholder(slide, x, y, w, h, label, sub="", icon="📷"):
    box = fill_rect(slide, x, y, w, h, LIGHT, line=PHBORD, line_w=Pt(1.5))
    # dashed border
    ln = box.line._get_or_add_ln()
    d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(d)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = icon
    r.font.size = Pt(34); r.font.name = FONT; r.font.color.rgb = TEAL
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(15); r2.font.bold = True; r2.font.name = FONT
    r2.font.color.rgb = NAVY
    if sub:
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run(); r3.text = sub
        r3.font.size = Pt(11); r3.font.italic = True; r3.font.name = FONT
        r3.font.color.rgb = GRAY
    return box


def chip(slide, x, y, text, color=RED, w=Inches(2.4), h=Inches(0.5)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True; r.font.name = FONT
    r.font.color.rgb = WHITE
    return shp


N = [0]
def page(slide):
    N[0] += 1
    footer(slide, N[0])


# =====================================================================
# 1. TITLE
# =====================================================================
s = add_slide()
fill_rect(s, 0, 0, SW, SH, NAVY)
fill_rect(s, 0, Inches(4.95), SW, Inches(0.10), TEAL)
fill_rect(s, 0, Inches(5.05), SW, Inches(2.45), WHITE)
txt(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(0.5),
    "SUMMER DISCOVERY  •  2026", size=18, color=TEAL, bold=True)
txt(s, Inches(0.85), Inches(1.7), Inches(11.7), Inches(2.2),
    ["Camp Health & Safety", "Counselor Training"], size=52, color=WHITE, bold=True,
    line_spacing=1.0)
txt(s, Inches(0.9), Inches(4.0), Inches(11), Inches(0.7),
    "Everything you need to keep our campers safe — and stay calm doing it.",
    size=19, color=RGBColor(0xCB, 0xE3, 0xDF), italic=True)
txt(s, Inches(0.9), Inches(5.45), Inches(11), Inches(1.3),
    [{"t": "Presented by the Camp Health Office", "size": 18, "bold": True, "color": NAVY},
     {"t": "Based on the Summer 2026 Staff Handbook, Section 9: Health, Safety & Wellness", "size": 14, "color": GRAY},
     {"t": "Your School Nurse  •  When in doubt, call. Always.", "size": 14, "color": TEAL, "bold": True}],
    space_after=6)
# no footer on title

# =====================================================================
# 2. WELCOME / WHY THIS MATTERS
# =====================================================================
s = add_slide()
header(s, "Welcome", "First Job? First Aid? You've Got This.")
bullets(s, Inches(0.7), Inches(1.7), Inches(7.0), Inches(5),
        ["You are the first set of eyes on every camper, every day.",
         "Most of what you'll handle is small: scrapes, sniffles, homesickness.",
         "A few things are serious. We'll make sure you know the difference.",
         "Golden rule: you are never expected to handle an emergency alone.",
         "When something feels off, you call the Health Office. No bravery points for guessing."],
        size=19, space_after=16)
ph = placeholder(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.2),
            "INSERT PHOTO", "Camp staff / smiling counselors (optional)")
chip(s, Inches(8.1), Inches(5.3), "Today's promise: keep it simple, keep it safe.",
     color=TEAL, w=Inches(4.6), h=Inches(0.6))
page(s)

# =====================================================================
# 3. AGENDA
# =====================================================================
s = add_slide()
header(s, "Roadmap", "What We'll Cover Today")
col1 = ["Your medical responsibilities", "Student illness & reporting",
        "Minor first aid & the first aid kit", "Basic first aid skills",
        "Hands-only CPR basics", "Doctor visits & paperwork"]
col2 = ["Medication management", "Mental health & wellness",
        "Allergies & anaphylaxis", "EpiPen, seizures & inhalers",
        "Other medical emergencies", "Map & route to the ER"]
bullets(s, Inches(0.8), Inches(1.9), Inches(5.8), Inches(5), col1, size=19, space_after=15)
bullets(s, Inches(7.0), Inches(1.9), Inches(5.8), Inches(5), col2, size=19, space_after=15)
page(s)

# =====================================================================
# 4. MEDICAL RESPONSIBILITIES (9.1)
# =====================================================================
s = add_slide()
header(s, "Section 9.1", "Your Medical Responsibilities")
bullets(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5),
        ["Take all symptoms and complaints of illness seriously — even the dramatic ones.",
         "Notify the Site Nurse or Wellness Coordinator immediately when something's wrong.",
         "Document all medical incidents according to protocol (if it's not written down, it didn't happen).",
         "Maintain confidentiality of every student's medical information.",
         "Be aware of allergies and medical conditions of the students in your group.",
         "Know how to respond to common health emergencies — that's what today is for."],
        size=19, space_after=15)
page(s)

# =====================================================================
# 5. STUDENT ILLNESS (9.1.1)
# =====================================================================
s = add_slide()
header(s, "Section 9.1.1", "When a Camper Feels Sick")
bullets(s, Inches(0.7), Inches(1.7), Inches(7.1), Inches(5),
        ["Believe them first, investigate second.",
         "Notify the Site Nurse / Wellness Coordinator right away.",
         "Document the incident per protocol.",
         "Keep their medical information private.",
         "Know your group's allergies and conditions before symptoms show up."],
        size=18, space_after=13)
fill_rect(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.6), LIGHT)
txt(s, Inches(8.35), Inches(2.05), Inches(4.2), Inches(0.5),
    "WATCH FOR", size=14, color=RED, bold=True)
bullets(s, Inches(8.35), Inches(2.55), Inches(4.2), Inches(3),
        ["Fever, chills, or flushing", "Vomiting or stomach pain",
         "Rash you can't explain", "Unusual sleepiness or confusion",
         "\"I just don't feel right\""],
        size=15, bcolor=RED, space_after=9)
page(s)

# =====================================================================
# 6. MINOR FIRST AID & KIT (9.1.2)
# =====================================================================
s = add_slide()
header(s, "Section 9.1.2", "Minor First Aid & the First Aid Kit")
bullets(s, Inches(0.7), Inches(1.7), Inches(7.1), Inches(5),
        ["Administer minor first aid for small scrapes, cuts, and bruises on arms and legs.",
         "Notify the Site Nurse / Health Center Coordinator when you use any supplies.",
         "Document minor first aid in the Medical Log.",
         "Complete an Incident Report for any accident requiring first aid.",
         "Wear protective gloves whenever blood or bodily fluids are present.",
         "Dispose of medical waste properly."],
        size=17, space_after=11)
placeholder(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.4),
            "INSERT PHOTO", "Contents of the camp first aid kit")
chip(s, Inches(8.1), Inches(5.45), "Gloves on BEFORE you play hero.",
     color=AMBER, w=Inches(4.6), h=Inches(0.6))
page(s)

# =====================================================================
# 7. BASIC FIRST AID  (abrasion placeholder)
# =====================================================================
s = add_slide()
header(s, "Skills", "Basic First Aid: Scrapes, Cuts & Abrasions", bar=TEAL)
bullets(s, Inches(0.7), Inches(1.7), Inches(7.1), Inches(5),
        ["Gloves on, calm voice on.",
         "Rinse the wound with clean water to flush out dirt and debris.",
         "Apply gentle pressure with clean gauze to stop bleeding.",
         "Pat dry, apply antibiotic ointment, and cover with a bandage.",
         "Tell the camper they were very brave (it helps).",
         "Log it, and flag deep, gaping, or won't-stop-bleeding wounds to the Nurse."],
        size=17, space_after=11)
placeholder(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.6),
            "INSERT PHOTO", "Example of an abrasion / scrape")
chip(s, Inches(8.1), Inches(5.6), "When in doubt, escalate — don't improvise.",
     color=RED, w=Inches(4.6), h=Inches(0.6))
page(s)

# =====================================================================
# 8. BASIC CPR  (CPR placeholder)
# =====================================================================
s = add_slide()
header(s, "Skills", "Basic CPR: Hands-Only Lifesaving", bar=RED)
chip(s, Inches(0.7), Inches(1.55), "Call 911 FIRST  •  Send someone for the AED & Nurse",
     color=RED, w=Inches(6.6), h=Inches(0.55))
bullets(s, Inches(0.7), Inches(2.3), Inches(7.1), Inches(4.5),
        ["Check responsiveness — tap and shout. No response, not breathing? Act.",
         "Call 911 (or point at someone: \"YOU, call 911\").",
         "Push hard and fast in the center of the chest — about 2 inches deep.",
         "Aim for 100–120 pushes a minute (beat of \"Stayin' Alive\").",
         "Don't stop until help, an AED, or the person takes over.",
         "Use the AED as soon as it arrives — it talks you through it."],
        size=17, space_after=11)
placeholder(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.7),
            "INSERT PHOTO / DIAGRAM", "How to perform CPR (hand placement)")
chip(s, Inches(8.1), Inches(5.7), "Formal CPR/AED certification is required.",
     color=NAVY, w=Inches(4.6), h=Inches(0.6))
page(s)

# =====================================================================
# 9. DOCTOR VISITS (9.1.3)
# =====================================================================
s = add_slide()
header(s, "Section 9.1.3", "Doctor Visits & Off-Site Appointments")
col1 = ["Staff must accompany students to all outside medical appointments.",
        "Doctor duty rotates among all staff.",
        "The Wellness RC is assigned a van each day.",
        "Bring the student's Medical Form and SD Medical Report to every appointment.",
        "A second staff member rides along if you're taking only one student.",
        "Wait patiently with the student during exams and treatment."]
col2 = ["Have the doctor / nurse speak directly with parents when possible.",
        "Complete the SD Medical Report before leaving the office.",
        "Get a written diagnosis and suggested treatment plan.",
        "Return all paperwork to the Site Nurse / Health Center Coordinator.",
        "Non-life-threatening? Inform the Director before calling parents."]
bullets(s, Inches(0.65), Inches(1.75), Inches(6.1), Inches(5.2), col1, size=15.5, space_after=10)
bullets(s, Inches(6.95), Inches(1.75), Inches(6.1), Inches(5.2), col2, size=15.5, space_after=10)
page(s)

# =====================================================================
# 10. INTERNATIONAL STUDENT VISITS (9.1.4)
# =====================================================================
s = add_slide()
header(s, "Section 9.1.4", "International Student Medical Visits")
bullets(s, Inches(0.8), Inches(1.85), Inches(11.5), Inches(5),
        ["International students have an ID card attached to their Medical Form — keep them together.",
         "For in-network facilities, students pay only the deductible at point of service.",
         "For out-of-network facilities, students will be billed later.",
         "Return all payment receipts and medical paperwork to the Site Nurse.",
         "Always keep the medical vehicles fully fueled and ready to use."],
        size=19, space_after=16)
chip(s, Inches(0.8), Inches(5.6), "Receipts matter. Pockets are not a filing system.",
     color=TEAL, w=Inches(6.2), h=Inches(0.6))
page(s)

# =====================================================================
# 11. MEDICATION MANAGEMENT (9.1.5)
# =====================================================================
s = add_slide()
header(s, "Section 9.1.5", "Medication Management")
bullets(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
        ["ALL prescription and over-the-counter meds are stored in the Summer Discovery Health Center.",
         "Counselors do NOT administer medications to students. Full stop.",
         "The Site Nurse is responsible for medication distribution and documentation.",
         "Report any concerns about medication compliance or misuse immediately.",
         "No nurse on campus? The Program Director gives specific instructions."],
        size=18, space_after=14, bcolor=RED)
chip(s, Inches(0.8), Inches(5.7), "Not even an Advil. The answer is \"let's see the Nurse.\"",
     color=RED, w=Inches(7.4), h=Inches(0.6))
page(s)

# =====================================================================
# 12. MENTAL HEALTH (9.1.6)
# =====================================================================
s = add_slide()
header(s, "Section 9.1.6", "Mental Health & Wellness Support", bar=TEAL)
bullets(s, Inches(0.7), Inches(1.75), Inches(7.2), Inches(5),
        ["Be alert to signs of mental health concerns in students.",
         "Report any concerning behaviors to the administration.",
         "Maintain appropriate confidentiality while keeping students safe.",
         "Support students facing homesickness, anxiety, or other emotional challenges.",
         "Know the referral process for mental health support services."],
        size=18, space_after=13)
fill_rect(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.7), LIGHT)
txt(s, Inches(8.35), Inches(2.05), Inches(4.2), Inches(0.5),
    "A GOOD COUNSELOR...", size=14, color=TEAL, bold=True)
bullets(s, Inches(8.35), Inches(2.55), Inches(4.2), Inches(3),
        ["Listens more than they fix", "Takes \"I miss home\" seriously",
         "Notices the quiet kid", "Loops in a supervisor early",
         "Knows safety beats secrecy"],
        size=15, bcolor=TEAL, space_after=10)
page(s)

# =====================================================================
# 13. FOOD ALLERGIES & ANAPHYLAXIS (9.2.1)
# =====================================================================
s = add_slide()
header(s, "Section 9.2.1", "Food Allergies & Anaphylaxis", bar=RED)
bullets(s, Inches(0.7), Inches(1.7), Inches(7.1), Inches(5),
        ["Know the signs: hives, swelling, difficulty breathing.",
         "Common allergens: peanuts, tree nuts, shellfish, dairy, eggs.",
         "Anaphylaxis is severe and potentially life-threatening.",
         "Symptoms: throat swelling, trouble breathing, drop in blood pressure, loss of consciousness.",
         "Know where the emergency meds (EpiPens) are kept.",
         "Know the emergency response steps cold."],
        size=16.5, space_after=10, bcolor=RED)
placeholder(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.6),
            "INSERT PHOTO", "Hives / allergic reaction example")
chip(s, Inches(8.1), Inches(5.6), "Seconds matter. Don't wait to \"see if it passes.\"",
     color=RED, w=Inches(4.6), h=Inches(0.6))
page(s)

# =====================================================================
# 14. EPIPEN ADMINISTRATION (9.2.2)
# =====================================================================
s = add_slide()
header(s, "Section 9.2.2", "EpiPen Administration — Step by Step", bar=RED)
steps = ["Remove the auto-injector from the carrier tube.",
         "Grasp it with the orange tip pointing downward.",
         "Remove the blue safety release (blue to the sky).",
         "Place the orange tip against the outer thigh at 90°.",
         "Push firmly until it \"clicks\" — hold for 10 seconds (orange to the thigh).",
         "Remove and massage the injection area.",
         "Call 911 immediately after administering.",
         "Contact the Site Nurse and program administration.",
         "Document the incident thoroughly.",
         "The student STILL goes to the hospital after an EpiPen."]
bullets(s, Inches(0.7), Inches(1.7), Inches(7.3), Inches(5.4),
        steps, size=15, space_after=7.5, bcolor=RED)
placeholder(s, Inches(8.2), Inches(1.8), Inches(4.5), Inches(3.5),
            "INSERT PHOTO", "EpiPen auto-injector / how to use")
chip(s, Inches(8.2), Inches(5.45), "Blue to the sky, orange to the thigh.",
     color=NAVY, w=Inches(4.5), h=Inches(0.6))
chip(s, Inches(8.2), Inches(6.15), "Always hospital after. No exceptions.",
     color=RED, w=Inches(4.5), h=Inches(0.6))
page(s)

# =====================================================================
# 15. SEIZURES  (placeholder)
# =====================================================================
s = add_slide()
header(s, "Emergency Skills", "Seizures: Stay Calm, Keep Them Safe", bar=RED)
bullets(s, Inches(0.7), Inches(1.7), Inches(7.1), Inches(5),
        ["Stay calm and note the time it starts.",
         "Ease the person to the floor and clear hard or sharp objects away.",
         "Turn them gently onto their side to keep the airway clear.",
         "Cushion their head; loosen anything tight around the neck.",
         "Do NOT hold them down and do NOT put anything in their mouth.",
         "Call 911 if it lasts over 5 minutes, repeats, or they don't wake up.",
         "Notify the Site Nurse and document everything."],
        size=16, space_after=9, bcolor=RED)
placeholder(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.6),
            "INSERT PHOTO / DIAGRAM", "Seizure first aid — recovery position")
chip(s, Inches(8.1), Inches(5.6), "Time it. Protect the head. Side position.",
     color=NAVY, w=Inches(4.6), h=Inches(0.6))
page(s)

# =====================================================================
# 16. ALBUTEROL INHALERS  (placeholder)
# =====================================================================
s = add_slide()
header(s, "Emergency Skills", "Albuterol Inhalers: Helping Campers Breathe", bar=TEAL)
bullets(s, Inches(0.7), Inches(1.7), Inches(7.1), Inches(5),
        ["Know which campers have asthma BEFORE the wheeze starts.",
         "Signs of trouble: wheezing, coughing, tight chest, hard time talking.",
         "The Nurse/plan directs inhaler use — follow the student's action plan.",
         "Use a spacer if one is provided; shake, one puff, slow deep breath, hold.",
         "Have them sit upright and stay calm; panic makes breathing harder.",
         "Not improving, lips/face turning blue, or no inhaler? Call 911 now.",
         "Notify the Site Nurse and document the episode."],
        size=15.5, space_after=8.5)
placeholder(s, Inches(8.1), Inches(1.85), Inches(4.6), Inches(3.6),
            "INSERT PHOTO", "Albuterol inhaler (+ spacer)")
chip(s, Inches(8.1), Inches(5.6), "Sit up, slow breaths, follow their plan.",
     color=TEAL, w=Inches(4.6), h=Inches(0.6))
page(s)

# =====================================================================
# 17. OTHER MEDICAL EMERGENCIES (9.2.3)
# =====================================================================
s = add_slide()
header(s, "Section 9.2.3", "Other Medical Emergencies", bar=RED)
bullets(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
        ["Stay calm and assess the situation.",
         "Call for medical assistance as appropriate — 911 for life-threatening emergencies.",
         "Never leave a student alone during a medical emergency.",
         "Contact the Site Nurse and program administration immediately.",
         "Document all actions taken.",
         "Follow up with the affected students and their parents."],
        size=19, space_after=14, bcolor=RED)
chip(s, Inches(0.8), Inches(6.0), "Calm is contagious. So is panic. Choose calm.",
     color=NAVY, w=Inches(6.4), h=Inches(0.6))
page(s)

# =====================================================================
# 18. MAP TO UCLA ER  (map placeholder)
# =====================================================================
s = add_slide()
header(s, "Know Before You Go", "Route to the Emergency Room — UCLA", bar=NAVY)
placeholder(s, Inches(0.7), Inches(1.75), Inches(7.6), Inches(4.9),
            "INSERT MAP", "Camp → Ronald Reagan UCLA Medical Center ER")
txt(s, Inches(8.5), Inches(1.8), Inches(4.3), Inches(0.5),
    "BEFORE YOU DRIVE", size=15, color=RED, bold=True)
bullets(s, Inches(8.5), Inches(2.3), Inches(4.3), Inches(4),
        ["Confirm the exact ER address & entrance",
         "Bring the camper's Medical Form + ID card",
         "Take a charged phone + the van keys",
         "Tell the Nurse & Director you're leaving",
         "Never transport a critical patient — call 911",
         "Save the ER + Health Office numbers in your phone"],
        size=14.5, bcolor=NAVY, space_after=11)
chip(s, Inches(8.5), Inches(6.15), "Life-threatening = 911, not the van.",
     color=RED, w=Inches(4.3), h=Inches(0.6))
page(s)

# =====================================================================
# 19. EMERGENCY QUICK REFERENCE
# =====================================================================
s = add_slide()
header(s, "Keep This Handy", "Emergency Quick Reference", bar=RED)
# three cards
cards = [
    ("CALL 911 IF...", RED,
     ["Not breathing / no pulse", "Severe allergic reaction",
      "Seizure over 5 min", "Uncontrolled bleeding", "Unconscious or confused"]),
    ("THEN, EVERY TIME", NAVY,
     ["Call the Site Nurse", "Notify the Director", "Stay with the student",
      "Write down what happened", "Follow up with parents"]),
    ("REMEMBER", TEAL,
     ["You're never alone in this", "Gloves before blood",
      "Counselors don't give meds", "Document everything",
      "When in doubt — escalate"]),
]
cx = Inches(0.7)
cw = Inches(3.95)
gap = Inches(0.18)
for title_t, col, items in cards:
    fill_rect(s, cx, Inches(1.75), cw, Inches(0.7), col)
    txt(s, cx, Inches(1.86), cw, Inches(0.5), title_t, size=16, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    fill_rect(s, cx, Inches(2.45), cw, Inches(4.0), LIGHT)
    bullets(s, cx + Inches(0.25), Inches(2.65), cw - Inches(0.5), Inches(3.6),
            items, size=15, bcolor=col, space_after=13)
    cx = Emu(cx + cw + gap)
page(s)

# =====================================================================
# 20. CLOSING
# =====================================================================
s = add_slide()
fill_rect(s, 0, 0, SW, SH, NAVY)
fill_rect(s, 0, Inches(2.6), SW, Inches(0.08), TEAL)
txt(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(1.4),
    ["You're ready.", "Now go be the calm one."], size=46, color=WHITE, bold=True)
txt(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(2),
    [{"t": "Three things to take with you:", "size": 20, "bold": True, "color": TEAL},
     {"t": "1.  Take every symptom seriously.", "size": 19, "color": WHITE},
     {"t": "2.  You are never expected to handle an emergency alone.", "size": 19, "color": WHITE},
     {"t": "3.  When in doubt, call the Health Office. Every single time.", "size": 19, "color": WHITE}],
    space_after=10)
txt(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1),
    [{"t": "Questions? Your School Nurse is here all summer.", "size": 18, "italic": True,
      "color": RGBColor(0xCB, 0xE3, 0xDF)},
     {"t": "Health Office: ____________   •   Emergency: 911", "size": 16, "bold": True,
      "color": WHITE}], space_after=8)

prs.save("Camp_Health_and_Safety_2026.pptx")
print("Saved Camp_Health_and_Safety_2026.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
